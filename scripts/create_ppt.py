from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# 색상 정의
BLACK = RGBColor(0, 0, 0)
DARK_BG = RGBColor(15, 15, 15)
WHITE = RGBColor(255, 255, 255)
ACCENT_YELLOW = RGBColor(255, 255, 0)
ACCENT_RED = RGBColor(255, 107, 107)
ACCENT_CYAN = RGBColor(78, 205, 196)
ACCENT_GOLD = RGBColor(255, 230, 109)
GRAY = RGBColor(128, 128, 128)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_GRAY = RGBColor(200, 200, 200)

def add_gradient_background(slide):
    """화려한 그라디언트 배경 추가"""
    # 메인 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # 좌상단 장식 원 (더 어둡게)
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(80, 30, 30)
    circle1.line.fill.background()

    # 우하단 장식 원 (더 어둡게)
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(25, 60, 60)
    circle2.line.fill.background()

    # 중앙 하단 악센트 (더 어둡게, 위치 조정)
    circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(5.5), Inches(3), Inches(3))
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = RGBColor(60, 60, 20)
    circle3.line.fill.background()

def add_simple_dark_bg(slide):
    """심플한 다크 배경"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

def add_accent_line(slide, x, y, width, color):
    """악센트 라인 추가"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_decorative_elements(slide):
    """장식 요소 추가"""
    # 우측 상단 작은 사각형들
    for i in range(3):
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(11.5 + i * 0.4), Inches(0.5 + i * 0.3),
            Inches(0.15), Inches(0.15)
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = ACCENT_YELLOW if i == 0 else (ACCENT_RED if i == 1 else ACCENT_CYAN)
        sq.line.fill.background()

    # 좌측 하단 라인
    for i in range(2):
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(6.5 + i * 0.2),
            Inches(1.5 - i * 0.5), Pt(2)
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = GRAY
        ln.line.fill.background()

def add_number_badge(slide, number, x, y, color):
    """숫자 배지 추가"""
    # 배경 원
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    # 숫자
    num_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.12), Inches(0.6), Inches(0.4))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER

# === 슬라이드 1: 표지 ===
def create_cover_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # 대형 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(12), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "iRUSH"
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DIGITAL CREATIVE AGENCY"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_YELLOW
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

    # 하단 설명
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(12), Inches(1))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Homepage Concept Presentation"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # 장식 라인
    add_accent_line(slide, 5.5, 5.2, 2.3, ACCENT_YELLOW)

    # 연도 표시
    year_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5))
    tf = year_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2024"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

# === 슬라이드 2: 프로젝트 개요 ===
def create_overview_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)
    add_decorative_elements(slide)

    # 섹션 번호
    add_number_badge(slide, "01", 0.8, 0.5, ACCENT_YELLOW)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT OVERVIEW"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_YELLOW)

    # 좌측 컬럼 - 프로젝트 정보
    info_items = [
        ("PROJECT", "iRUSH 홈페이지 리뉴얼"),
        ("CLIENT", "iRUSH"),
        ("TYPE", "Corporate Website"),
        ("YEAR", "2024")
    ]

    y_pos = 1.8
    for label, value in info_items:
        # 라벨
        lbl_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(2), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = ACCENT_YELLOW
        p.font.bold = True

        # 값
        val_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos + 0.35), Inches(4), Inches(0.5))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE

        y_pos += 1

    # 우측 컬럼 - 목적 및 타겟
    # 박스 배경
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_GRAY
    box.line.color.rgb = RGBColor(60, 60, 60)
    box.line.width = Pt(1)

    # 목적
    obj_title = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5), Inches(0.5))
    tf = obj_title.text_frame
    p = tf.paragraphs[0]
    p.text = "OBJECTIVE"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True

    obj_text = slide.shapes.add_textbox(Inches(7), Inches(2.6), Inches(5.2), Inches(1.5))
    tf = obj_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "24년간의 경험과 전문성을 현대적인 디자인으로 표현하고, 자동화된 IT 트렌드 뉴스 서비스로 업계 선도 기업 이미지 강화"
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_GRAY

    # 핵심 메시지
    msg_title = slide.shapes.add_textbox(Inches(7), Inches(4.3), Inches(5), Inches(0.5))
    tf = msg_title.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY MESSAGE"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True

    msg_text = slide.shapes.add_textbox(Inches(7), Inches(4.8), Inches(5.2), Inches(1))
    tf = msg_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FEVER to Create • BLISS to Deliver • FAITH to Succeed"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True

# === 슬라이드 3: 핵심 가치 ===
def create_values_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)

    # 섹션 번호
    add_number_badge(slide, "02", 0.8, 0.5, ACCENT_RED)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CORE VALUES"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_RED)

    values = [
        ("FEVER", "to Create", "창작에 대한 뜨거운 열정.\n매 프로젝트에 혼신의 열정을\n쏟아붓습니다.", ACCENT_RED),
        ("BLISS", "to Deliver", "전달하는 기쁨.\n완벽한 결과물을 전달할 때\n느끼는 충만한 행복.", ACCENT_CYAN),
        ("FAITH", "to Succeed", "성공에 대한 믿음.\n고객과 함께 성공할 수 있다는\n굳건한 신념.", ACCENT_GOLD)
    ]

    start_x = 1
    for i, (title, subtitle, desc, color) in enumerate(values):
        x = start_x + i * 4

        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.8), Inches(3.8), Inches(5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(25, 25, 25)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # 상단 악센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.8), Inches(3.8), Inches(0.15)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = color
        accent_bar.line.fill.background()

        # 번호
        num_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.2), Inches(1), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(14)
        p.font.color.rgb = color

        # 타이틀
        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.7), Inches(3.2), Inches(1))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = color

        # 서브타이틀
        s_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.8), Inches(3.2), Inches(0.5))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.italic = True

        # 설명
        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(4.5), Inches(3.2), Inches(2))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.line_spacing = 1.5

# === 슬라이드 4: 통계 ===
def create_stats_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)

    # 큰 배경 텍스트
    bg_text = slide.shapes.add_textbox(Inches(-1), Inches(1.5), Inches(15), Inches(5))
    tf = bg_text.text_frame
    p = tf.paragraphs[0]
    p.text = "STATISTICS"
    p.font.size = Pt(150)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)
    p.alignment = PP_ALIGN.CENTER

    # 섹션 번호
    add_number_badge(slide, "03", 0.8, 0.5, ACCENT_CYAN)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY STATISTICS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    stats = [
        ("24", "+", "YEARS", "2002년부터 시작된 여정", ACCENT_YELLOW),
        ("500", "+", "PROJECTS", "다양한 프로젝트 경험", ACCENT_RED),
        ("100", "%", "PASSION", "고객 만족 최우선", ACCENT_CYAN)
    ]

    start_x = 1.2
    for i, (num, suffix, label, desc, color) in enumerate(stats):
        x = start_x + i * 4

        # 숫자
        num_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.5), Inches(2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = num
        run.font.size = Pt(100)
        run.font.bold = True
        run.font.color.rgb = color

        # suffix
        run2 = p.add_run()
        run2.text = suffix
        run2.font.size = Pt(50)
        run2.font.bold = True
        run2.font.color.rgb = WHITE

        # 라벨
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(4.7), Inches(3.5), Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 구분선
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.3), Inches(2), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.fill.background()

        # 설명
        desc_box = slide.shapes.add_textbox(Inches(x), Inches(5.5), Inches(3.5), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

# === 슬라이드 5: 기술 스택 ===
def create_tech_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)
    add_decorative_elements(slide)

    # 섹션 번호
    add_number_badge(slide, "04", 0.8, 0.5, ACCENT_GOLD)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TECH STACK"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_GOLD)

    tech_items = [
        ("FRONTEND", [
            ("React 19", "최신 UI 라이브러리"),
            ("Vite", "초고속 빌드 도구"),
            ("GSAP", "프리미엄 애니메이션"),
            ("Framer Motion", "부드러운 트랜지션")
        ], ACCENT_RED),
        ("BACKEND", [
            ("GitHub Actions", "자동화 워크플로우"),
            ("Cheerio", "웹 크롤링 엔진"),
            ("Node.js", "서버 스크립트")
        ], ACCENT_CYAN),
        ("DEPLOY", [
            ("Vercel", "자동 배포"),
            ("GitHub", "버전 관리"),
            ("CDN", "글로벌 배포")
        ], ACCENT_GOLD)
    ]

    start_x = 0.8
    for i, (category, items, color) in enumerate(tech_items):
        x = start_x + i * 4.2

        # 카테고리 헤더
        cat_box = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(3.8), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # 아이템들
        y = 2.3
        for name, desc in items:
            # 배경 박스
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(3.8), Inches(1)
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = DARK_GRAY
            item_bg.line.fill.background()

            # 이름
            name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(3.4), Inches(0.4))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # 설명
            desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.55), Inches(3.4), Inches(0.4))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = GRAY

            y += 1.15

# === 슬라이드 6: 디자인 컨셉 ===
def create_design_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)

    # 섹션 번호
    add_number_badge(slide, "05", 0.8, 0.5, ACCENT_RED)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DESIGN CONCEPT"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_RED)

    # 컬러 팔레트
    colors_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(3), Inches(0.5))
    tf = colors_title.text_frame
    p = tf.paragraphs[0]
    p.text = "COLOR PALETTE"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_YELLOW
    p.font.bold = True

    palette = [
        (DARK_BG, "Dark BG", "#0F0F0F"),
        (ACCENT_YELLOW, "Accent", "#FFFF00"),
        (ACCENT_RED, "Energy", "#FF6B6B"),
        (ACCENT_CYAN, "Trust", "#4ECDC4"),
        (WHITE, "Text", "#FFFFFF")
    ]

    for i, (color, name, hex_code) in enumerate(palette):
        x = 0.8 + i * 1.3
        # 컬러 박스
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(1.1), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(60, 60, 60)
        box.line.width = Pt(1)

        # 이름
        n_box = slide.shapes.add_textbox(Inches(x), Inches(3.4), Inches(1.1), Inches(0.3))
        tf = n_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # HEX
        h_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(1.1), Inches(0.3))
        tf = h_box.text_frame
        p = tf.paragraphs[0]
        p.text = hex_code
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # 타이포그래피
    typo_title = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(5), Inches(0.5))
    tf = typo_title.text_frame
    p = tf.paragraphs[0]
    p.text = "TYPOGRAPHY"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_CYAN
    p.font.bold = True

    fonts = [
        ("Bebas Neue", "DISPLAY", "대형 타이틀용"),
        ("Inter", "BODY", "본문 가독성"),
        ("Noto Sans KR", "KOREAN", "한글 최적화")
    ]

    y = 2.2
    for font, role, desc in fonts:
        # 폰트 박스
        f_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(y), Inches(5), Inches(0.9))
        f_box.fill.solid()
        f_box.fill.fore_color.rgb = DARK_GRAY
        f_box.line.fill.background()

        # 폰트명
        fn_box = slide.shapes.add_textbox(Inches(7.7), Inches(y + 0.1), Inches(3), Inches(0.4))
        tf = fn_box.text_frame
        p = tf.paragraphs[0]
        p.text = font
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 역할
        r_box = slide.shapes.add_textbox(Inches(7.7), Inches(y + 0.5), Inches(2), Inches(0.3))
        tf = r_box.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(10)
        p.font.color.rgb = ACCENT_YELLOW

        # 설명
        d_box = slide.shapes.add_textbox(Inches(10), Inches(y + 0.3), Inches(2.3), Inches(0.4))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

        y += 1.05

    # 비주얼 스타일
    style_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(5), Inches(0.5))
    tf = style_title.text_frame
    p = tf.paragraphs[0]
    p.text = "VISUAL STYLE"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True

    styles = [
        "다크 테마 기반의 프리미엄 디자인",
        "대형 타이포그래피와 충분한 여백",
        "스크롤 트리거 애니메이션",
        "인터랙티브 파티클 효과"
    ]

    y = 4.8
    for style in styles:
        s_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(6), Inches(0.4))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"→  {style}"
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        y += 0.45

# === 슬라이드 7: 사이트 구조 ===
def create_structure_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)

    # 섹션 번호
    add_number_badge(slide, "06", 0.8, 0.5, ACCENT_CYAN)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SITE STRUCTURE"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_CYAN)

    pages = [
        ("MAIN", "메인 페이지", ["히어로 비디오 섹션", "핵심 가치 소개", "서비스 영역", "프로젝트 미리보기"], ACCENT_YELLOW),
        ("ABOUT", "회사 소개", ["Since 2002 히스토리", "철학 및 비전", "팀 구성", "클라이언트"], ACCENT_RED),
        ("WORK", "프로젝트", ["Featured 슬라이드쇼", "프로젝트 아카이브", "스와이프 네비게이션"], ACCENT_CYAN),
        ("TREND", "트렌드 데스크", ["자동 수집 IT 뉴스", "파티클 배경", "아코디언 UI"], ACCENT_GOLD)
    ]

    start_x = 0.6
    for i, (name, title, features, color) in enumerate(pages):
        x = start_x + i * 3.2

        # 페이지 카드
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.7), Inches(3), Inches(5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(25, 25, 25)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 페이지 이름
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.9), Inches(2.6), Inches(0.6))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color

        # 한글 타이틀
        kr_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.5), Inches(2.6), Inches(0.4))
        tf = kr_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE

        # 구분선
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.2), Inches(3), Inches(2.6), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(60, 60, 60)
        line.line.fill.background()

        # 기능 목록
        y = 3.2
        for feature in features:
            f_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(2.6), Inches(0.5))
            tf = f_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {feature}"
            p.font.size = Pt(11)
            p.font.color.rgb = GRAY
            y += 0.5

# === 슬라이드 8: 주요 기능 ===
def create_features_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)
    add_decorative_elements(slide)

    # 섹션 번호
    add_number_badge(slide, "07", 0.8, 0.5, ACCENT_YELLOW)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY FEATURES"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_YELLOW)

    features = [
        ("AUTO CRAWLING", "자동 뉴스 크롤링", "매일 오전 9시 GitHub Actions로\nZDNet Korea, ITWorld Korea에서\n최신 IT 뉴스를 자동 수집", ACCENT_RED),
        ("ANIMATIONS", "인터랙티브 애니메이션", "GSAP & Framer Motion 기반\n스크롤 트리거, 패럴랙스,\n마우스 반응형 파티클", ACCENT_CYAN),
        ("RESPONSIVE", "반응형 디자인", "모바일, 태블릿, 데스크톱\n모든 디바이스에 최적화된\n사용자 경험 제공", ACCENT_GOLD)
    ]

    start_x = 0.8
    for i, (title, subtitle, desc, color) in enumerate(features):
        x = start_x + i * 4.2

        # 아이콘 영역
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(1.8), Inches(1), Inches(1))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = color
        icon_bg.line.fill.background()

        # 아이콘 텍스트
        icon_text = slide.shapes.add_textbox(Inches(x + 1.3), Inches(2.05), Inches(1), Inches(0.6))
        tf = icon_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"0{i+1}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        # 타이틀
        t_box = slide.shapes.add_textbox(Inches(x), Inches(3), Inches(3.8), Inches(0.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 서브타이틀
        s_box = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(3.8), Inches(0.4))
        tf = s_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 설명 박스
        desc_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(4.1), Inches(3.8), Inches(2.2)
        )
        desc_bg.fill.solid()
        desc_bg.fill.fore_color.rgb = DARK_GRAY
        desc_bg.line.fill.background()

        # 설명
        d_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4.3), Inches(3.4), Inches(1.8))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.line_spacing = 1.4

# === 슬라이드 9: 클라이언트 ===
def create_clients_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_simple_dark_bg(slide)

    # 섹션 번호
    add_number_badge(slide, "08", 0.8, 0.5, ACCENT_RED)

    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MAJOR CLIENTS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_accent_line(slide, 1.6, 1.2, 2, ACCENT_RED)

    clients = [
        ("대기업", ["LG전자", "LG헬로비전", "삼성전자", "삼성SDS", "SK텔레콤", "현대자동차"], ACCENT_YELLOW),
        ("금융", ["신한은행", "신한금융지주", "KB국민은행"], ACCENT_CYAN),
        ("IT / 플랫폼", ["카카오", "네이버", "CJ ENM", "CJ올리브네트웍스"], ACCENT_RED),
        ("기타", ["국립암센터", "삼성문화재단", "두산", "하나제약"], ACCENT_GOLD)
    ]

    start_x = 0.8
    for i, (category, names, color) in enumerate(clients):
        x = start_x + i * 3.2

        # 카테고리 헤더
        cat_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.7), Inches(3), Inches(0.6)
        )
        cat_bg.fill.solid()
        cat_bg.fill.fore_color.rgb = color
        cat_bg.line.fill.background()

        cat_box = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(3), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER

        # 클라이언트 목록
        y = 2.5
        for name in names:
            n_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(2.6), Inches(0.5))
            tf = n_box.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            y += 0.55

# === 슬라이드 10: 마무리 ===
def create_ending_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # 대형 Thank You
    thanks_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(12), Inches(2))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(90)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 회사명
    company_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(12), Inches(0.8))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "iRUSH - Digital Creative Agency"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_YELLOW
    p.alignment = PP_ALIGN.CENTER

    # 연락처
    contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(12), Inches(0.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "info@irush.co.kr"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # 하단 핵심 가치
    values_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(12), Inches(0.5))
    tf = values_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FEVER  •  BLISS  •  FAITH"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# === 슬라이드 생성 실행 ===
create_cover_slide()
create_overview_slide()
create_values_slide()
create_stats_slide()
create_tech_slide()
create_design_slide()
create_structure_slide()
create_features_slide()
create_clients_slide()
create_ending_slide()

# PPT 저장
output_path = os.path.join(os.path.dirname(__file__), '..', 'iRUSH_Homepage_Concept.pptx')
prs.save(output_path)
print(f"PPT 생성 완료: {os.path.abspath(output_path)}")
